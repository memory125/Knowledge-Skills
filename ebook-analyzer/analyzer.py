#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Ebook Analyzer v2.0 - 电子书智能分析工具 (升级版)
完全本地化，无需外部 API，支持 37 维深度分析
"""

import os
import sys
import json
import argparse
import re
from pathlib import Path
from typing import Dict, List, Optional, Set
from datetime import datetime

# PDF 解析
try:
    import pdfplumber
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

# Word 文档解析
try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# PPT 解析
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# EPUB 解析
try:
    import ebooklib
    from ebooklib import epub
    HAS_EPUB = True
except ImportError:
    HAS_EPUB = False


# ==================== 分析维度定义 ====================

DIMENSION_CATEGORIES = {
    "core": {
        "1_内容概述": ("📖 本书内容", "300-500 字详细概述全书的主要内容、章节结构和核心议题"),
        "2_一句话总结": ("💡 一句话总结", "用一句精炼的话概括全书核心，不超过 50 字"),
        "3_结构图谱": ("🗺️ 逻辑结构图谱", "如果是小说：梳理人物关系图；如果是非虚构：绘制概念关联图"),
        "4_本书要点": ("🔑 本书要点", "提取 5-10 个核心观点/情节要点"),
        "5_金句汇总": ("✨ 金句汇总", "摘录 5-8 句最有价值的名言警句，标注出处"),
        "6_作者思想": ("🎯 作者思想", "分析作者的核心论点、价值观和世界观"),
        "7_现实意义": ("🌍 现实意义", "这本书对当下读者的启发和实用价值"),
        "8_适合读者": ("👥 适合读者", "描述最适合的读者群体"),
        "9_书籍评分": ("⭐ 书籍评分", "从内容深度、可读性、启发性三个维度打分并说明理由"),
    },
    "structure": {
        "10_写作手法": ("✍️ 写作手法分析", "叙事技巧、修辞手法、语言风格、结构安排"),
        "11_章节脉络": ("📚 章节脉络图", "每章核心任务、章节间逻辑连接、情节/论点的递进关系"),
        "12_关键词云": ("☁️ 关键词云", "提取全书最高频的 50-100 个词，生成概念统计"),
    },
    "content": {
        "13_核心问题": ("❓ 核心问题清单", "列出本书试图回答的 3-5 个根本问题"),
        "14_概念词典": ("📖 概念词典", "提取并定义本书的 10-20 个核心术语"),
        "15_论证逻辑链": ("🔗 论证逻辑链", "分析前提、证据、推理、结论的严密性（仅非虚构）"),
        "16_人物弧光": ("👤 人物弧光分析", "追踪主要人物的成长轨迹（仅小说/传记）"),
        "17_冲突类型": ("⚔️ 冲突类型分析", "分析内部冲突、外部冲突、关系冲突等（仅小说）"),
    },
    "knowledge": {
        "18_知识图谱": ("🌐 知识图谱", "绘制本书涉及的知识网络和概念关联"),
        "19_方法论清单": ("🛠️ 方法论清单", "提取可实操的方法步骤"),
        "20_案例库": ("📋 案例库", "摘录本书中的典型案例及其启示"),
    },
    "critical": {
        "21_潜在偏见": ("⚖️ 潜在偏见识别", "分析作者的时代局限、文化偏见、利益关联等"),
        "22_逻辑漏洞": ("🔍 逻辑漏洞检测", "发现概念偷换、因果倒置、以偏概全等问题"),
        "23_反方观点": ("↔️ 反方观点", "思考如果反对这本书的观点，会怎么说"),
        "24_争议点整理": ("🎭 争议点整理", "列出学界/读者对本书的主要争议"),
    },
    "practice": {
        "25_行动清单": ("✅ 行动清单", "将知识转化为具体可执行的行动"),
        "26_习惯计划": ("📅 习惯养成计划", "基于本书设计 21/30 天改变计划"),
        "27_讨论问题": ("💬 讨论问题集", "为读书会准备的 10-15 个深入问题"),
        "28_复习要点": ("📝 复习要点", "制作便于记忆的知识卡片"),
    },
    "cross": {
        "29_同类对比": ("🔀 同类书对比", "与同领域其他经典进行多维度对比"),
        "30_思想渊源": ("🌳 思想渊源追溯", "分析本书的思想来源和知识传承"),
        "31_后续影响": ("📢 后续影响评估", "这本书影响了谁，引发了什么讨论"),
    },
    "growth": {
        "32_共鸣记录": ("❤️ 共鸣记录", "哪些内容最触动你，为什么"),
        "33_认知升级": ("🧠 认知升级点", "记录本书打破或更新了哪些旧观念"),
        "34_创作素材": ("✏️ 引用创作素材库", "为写作积累的典故、比喻、论据、故事"),
    },
    "meta": {
        "35_阅读反思": ("📔 阅读反思日志", "记录阅读过程中的思维变化"),
        "36_知识定位": ("🗺️ 知识体系定位", "绘制这本书在个人知识结构中的位置"),
        "37_未解问题": ("❔ 未解答的问题", "记录读完后留下的疑问和下一步学习方向"),
    }
}

# 预设模式配置
PRESET_MODES = {
    "quick": {
        "name": "快速筛选",
        "dimensions": ["1_内容概述", "2_一句话总结", "4_本书要点", "8_适合读者"],
        "description": "快速了解全书核心内容和价值"
    },
    "core": {
        "name": "核心九维",
        "dimensions": list(DIMENSION_CATEGORIES["core"].keys()),
        "description": "经典的九个基础分析维度"
    },
    "deep": {
        "name": "深度阅读",
        "dimensions": (
            list(DIMENSION_CATEGORIES["core"].keys()) +
            ["10_写作手法", "11_章节脉络", "12_关键词云"] +
            ["13_核心问题", "14_概念词典"] +
            ["18_知识图谱", "19_方法论清单", "20_案例库"] +
            ["32_共鸣记录", "33_认知升级"]
        ),
        "description": "深度学习和理解，适合精读"
    },
    "academic": {
        "name": "学术研究",
        "dimensions": (
            ["1_内容概述", "2_一句话总结", "6_作者思想"] +
            ["13_核心问题", "14_概念词典", "15_论证逻辑链", "18_知识图谱"] +
            ["21_潜在偏见", "22_逻辑漏洞", "23_反方观点", "24_争议点整理"] +
            ["29_同类对比", "30_思想渊源", "31_后续影响"] +
            ["37_未解问题"]
        ),
        "description": "批判性分析和学术引用"
    },
    "skill": {
        "name": "技能提升",
        "dimensions": (
            ["1_内容概述", "2_一句话总结", "5_金句汇总"] +
            ["19_方法论清单", "20_案例库"] +
            ["25_行动清单", "26_习惯计划", "28_复习要点"]
        ),
        "description": "提取方法和行动指南，快速应用"
    },
    "literature": {
        "name": "文学欣赏",
        "dimensions": (
            ["1_内容概述", "2_一句话总结", "5_金句汇总"] +
            ["10_写作手法", "11_章节脉络", "16_人物弧光", "17_冲突类型"] +
            ["34_创作素材"]
        ),
        "description": "分析文学技巧和艺术特色"
    },
    "discussion": {
        "name": "读书分享",
        "dimensions": (
            ["1_内容概述", "2_一句话总结", "6_作者思想", "7_现实意义"] +
            ["27_讨论问题", "32_共鸣记录", "33_认知升级"]
        ),
        "description": "为读书会和分享做准备"
    },
    "all": {
        "name": "全维度分析",
        "dimensions": [],  # 将在运行时填充所有维度
        "description": "完整的 37 维深度分析（推荐长书使用）"
    }
}


class EbookAnalyzer:
    """电子书分析器 v2.0"""
    
    def __init__(self, book_path: str):
        self.book_path = Path(book_path)
        self.text_content = ""
        self.metadata = {}
        self.sections = []
        self.word_freq = {}
        
    def extract_text(self) -> str:
        """提取文本内容"""
        ext = self.book_path.suffix.lower()
        
        if ext == '.pdf':
            return self._extract_pdf()
        elif ext in ['.docx', '.doc']:
            return self._extract_docx()
        elif ext in ['.pptx', '.ppt']:
            return self._extract_pptx()
        elif ext == '.epub':
            return self._extract_epub()
        else:
            try:
                with open(self.book_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except:
                return ""
    
    def _extract_pdf(self) -> str:
        """提取 PDF 内容"""
        if not HAS_PDF:
            print("错误：未安装 pdfplumber，请运行：pip install --break-system-packages pdfplumber")
            return ""
        
        text_parts = []
        try:
            with pdfplumber.open(self.book_path) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text:
                        self.metadata['page_count'] = len(pdf.pages)
                        text_parts.append(f"[第{i}页]\n{text.strip()}")
        except Exception as e:
            print(f"PDF 提取错误：{e}")
        
        return "\n\n".join(text_parts)
    
    def _extract_docx(self) -> str:
        """提取 DOCX 内容"""
        if not HAS_DOCX:
            print("错误：未安装 python-docx，请运行：pip install --break-system-packages python-docx")
            return ""
        
        try:
            doc = Document(self.book_path)
            self.metadata['title'] = doc.core_properties.title or self.book_path.stem
            self.metadata['author'] = doc.core_properties.author or "未知作者"
            
            text_parts = []
            current_section = []
            
            for paragraph in doc.paragraphs:
                text = paragraph.text.strip()
                if text:
                    # 检测标题（简单判断：无句号、较短）
                    if not any(c in text for c in '。！？!?.') and len(text) < 30:
                        if current_section:
                            text_parts.append('\n\n'.join(current_section))
                            current_section = []
                        text_parts.append(f"## {text}")
                    else:
                        current_section.append(text)
            
            if current_section:
                text_parts.append('\n\n'.join(current_section))
                
            return "\n\n".join(text_parts)
        except Exception as e:
            print(f"DOCX 提取错误：{e}")
            return ""
    
    def _extract_pptx(self) -> str:
        """提取 PPTX 内容"""
        if not HAS_PPTX:
            print("错误：未安装 python-pptx，请运行：pip install --break-system-packages python-pptx")
            return ""
        
        try:
            prs = Presentation(self.book_path)
            text_parts = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_items = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_items.append(shape.text.strip())
                
                if slide_items:
                    text_parts.append(f"[幻灯片 {i}]\n" + "\n".join(slide_items))
            
            return "\n\n".join(text_parts)
        except Exception as e:
            print(f"PPTX 提取错误：{e}")
            return ""
    
    def _extract_epub(self) -> str:
        """提取 EPUB 内容"""
        if not HAS_EPUB:
            print("错误：未安装 ebooklib，请运行：pip install --break-system-packages ebooklib")
            return ""
        
        try:
            book = epub.read_epub(self.book_path)
            text_parts = []
            
            # 提取元数据
            self.metadata['title'] = book.get_metadata('DC', 'title')
            self.metadata['author'] = book.get_metadata('DC', 'creator')
            
            for item in book.items:
                if hasattr(item, 'get_content') and item.get_content():
                    try:
                        content = item.get_content().decode('utf-8', errors='ignore')
                        # 简单清理 HTML 标签
                        content = re.sub(r'<[^>]+>', '', content)
                        text_parts.append(content.strip())
                    except:
                        pass
            
            return "\n\n".join(text_parts)
        except Exception as e:
            print(f"EPUB 提取错误：{e}")
            return ""
    
    def analyze_word_frequency(self, top_n=100) -> Dict[str, int]:
        """分析词频"""
        # 简单分词（中文按字符，英文按单词）
        if any('\u4e00' <= c <= '\u9fff' for c in self.text_content[:100]):
            # 中文：提取常用词
            words = re.findall(r'[\u4e00-\u9fff]{2,}', self.text_content)
        else:
            # 英文：提取单词
            words = re.findall(r'\b[a-zA-Z]{3,}\b', self.text_content.lower())
        
        # 统计词频
        freq = {}
        for word in words:
            if len(word) >= 2:  # 忽略单字/双字母
                freq[word] = freq.get(word, 0) + 1
        
        # 返回前 N 个高频词
        return dict(sorted(freq.items(), key=lambda x: x[1], reverse=True)[:top_n])
    
    def analyze_with_ollama(self, model: str = "qwen2.5", dimensions: List[str] = None) -> Dict:
        """使用 Ollama 进行智能分析"""
        import subprocess
        
        # 如果没有指定维度，使用核心九维
        if not dimensions:
            dimensions = list(DIMENSION_CATEGORIES["core"].keys())
        
        prompt = self._create_analysis_prompt(dimensions)
        
        try:
            result = subprocess.run(
                ['ollama', 'run', model, prompt],
                capture_output=True,
                text=True,
                timeout=600  # 10 分钟超时（深度分析可能需要更长时间）
            )
            
            if result.returncode == 0:
                return {
                    'success': True,
                    'analysis': result.stdout,
                    'dimensions_analyzed': dimensions
                }
            else:
                return {
                    'success': False,
                    'error': result.stderr,
                    'dimensions_analyzed': dimensions
                }
        except FileNotFoundError:
            return {
                'success': False,
                'error': '未找到 Ollama，请确保已安装并运行：ollama serve',
                'dimensions_analyzed': dimensions
            }
        except subprocess.TimeoutExpired:
            return {
                'success': False,
                'error': '分析超时（超过 10 分钟），可能是书籍内容过长或维度太多',
                'dimensions_analyzed': dimensions
            }
    
    def _create_analysis_prompt(self, dimensions: List[str]) -> str:
        """创建分析提示词"""
        
        # 提取书籍元信息
        title = self.metadata.get('title', self.book_path.stem)
        author = self.metadata.get('author', '未知')
        
        # 截取前一部分内容用于分析
        content_preview = self.text_content[:15000] if len(self.text_content) > 15000 else self.text_content
        
        # 构建维度描述
        dimension_descriptions = []
        for dim_key in dimensions:
            # 查找对应的类别和描述
            for category, dims in DIMENSION_CATEGORIES.items():
                if dim_key in dims:
                    icon, desc = dims[dim_key]
                    dimension_descriptions.append(f"- {icon} {dim_key.replace('_', ' ')}: {desc}")
                    break
        
        prompt = f"""# 电子书智能分析任务

## 📚 书籍基本信息
- **书名**: {title}
- **作者**: {author}
- **文件**: {self.book_path.name}
- **总字数**: {len(self.text_content):,} 字
- **分析日期**: {datetime.now().strftime('%Y-%m-%d %H:%M')}

## 📖 待分析的维度 ({len(dimensions)}个)
{'\n'.join(dimension_descriptions)}

---

## 📝 书籍内容（前 15000 字符预览）

{content_preview}

---

## ⚙️ 分析要求

### 通用标准
1. **基于文本**: 所有分析必须基于提供的文本内容，不要编造
2. **深度优先**: 每个维度都要有实质性内容，避免泛泛而谈
3. **证据支持**: 引用原文或具体情节/章节支撑观点
4. **结构清晰**: 使用标题、列表、表格等方式组织输出
5. **语言精炼**: 用简洁有力的中文表达

### 特殊要求
- **小说类**: 重点关注人物关系、情节发展、冲突类型
- **非虚构类**: 重点关注论点论证、概念定义、方法论
- **学术著作**: 重点关注理论框架、研究方法、知识贡献

---

## 📤 请输出完整的分析报告

请严格按照以下格式输出（使用 Markdown）：

# 《{title}》深度阅读分析

> 作者：{author}  
> 分析日期：{datetime.now().strftime('%Y-%m-%d')}  
> 分析维度：{len(dimensions)}维

---

"""
        
        # 添加每个维度的输出模板
        for dim_key in dimensions:
            for category, dims in DIMENSION_CATEGORIES.items():
                if dim_key in dims:
                    icon, _ = dims[dim_key]
                    title_safe = dim_key.replace('_', ' ')
                    prompt += f"## {icon} {title_safe}\n\n"
                    
                    # 根据维度类型添加具体指导
                    if "写作手法" in dim_key:
                        prompt += "【叙事技巧、修辞手法、语言风格、结构安排】\n"
                    elif "人物弧光" in dim_key or "冲突类型" in dim_key:
                        prompt += "【如文本内容不足，请说明并基于现有信息分析】\n"
                    elif "论证逻辑链" in dim_key:
                        prompt += "【前提 → 证据 → 推理 → 结论，评估严密性】\n"
                    elif "方法论清单" in dim_key or "行动清单" in dim_key:
                        prompt += "【步骤清晰、可操作、具体可行】\n"
                    elif "金句汇总" in dim_key:
                        prompt += "【5-8 句，标注出处，格式：> \"引文\" - 来源】\n"
                    elif "书籍评分" in dim_key:
                        prompt += "| 维度 | 评分 (1-5) | 理由 |\n|------|----------|------|\n"
                    else:
                        prompt += "【深入分析，提供具体内容】\n"
                    
                    prompt += "\n---\n\n"
                    break
        
        return prompt


def parse_dimensions(dim_str: str) -> List[str]:
    """解析维度字符串（如 "1,3,5-8,20"）"""
    dimensions = set()
    
    for part in dim_str.split(','):
        if '-' in part:  # 范围
            start, end = part.split('-')
            start_num, end_num = int(start), int(end)
            
            for cat_dims in DIMENSION_CATEGORIES.values():
                for key in cat_dims.keys():
                    key_num = int(key.split('_')[0])
                    if start_num <= key_num <= end_num:
                        dimensions.add(key)
        else:  # 单个维度
            num = int(part)
            for cat_dims in DIMENSION_CATEGORIES.values():
                for key in cat_dims.keys():
                    if int(key.split('_')[0]) == num:
                        dimensions.add(key)
    
    return sorted(list(dimensions), key=lambda x: int(x.split('_')[0]))


def main():
    parser = argparse.ArgumentParser(
        description='电子书智能分析工具 v2.0 - 37 维深度分析，完全本地化',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例：
  # 快速筛选模式
  ./ebook-analyzer analyze book.pdf --mode quick

  # 深度阅读模式（推荐）
  ./ebook-analyzer analyze book.epub --mode deep

  # 学术研究模式
  ./ebook-analyzer analyze thesis.pdf --mode academic

  # 自定义维度组合
  ./ebook-analyzer analyze book.docx --dimensions "1,2,5-9,19,25"

  # 导出为文件
  ./ebook-analyzer analyze book.pdf -m deep -o analysis.md

        """
    )
    
    parser.add_argument('command', choices=['analyze', 'extract'], 
                       help='命令：analyze(分析) 或 extract(提取文本)')
    parser.add_argument('book_path', type=str, help='书籍文件路径')
    parser.add_argument('--model', '-m', default='qwen2.5', 
                       help='Ollama 模型名称 (默认：qwen2.5)')
    parser.add_argument('--dimensions', '-d', type=str, default=None,
                       help='指定分析维度，如 "1,3,5-9" 或模式名 "deep"')
    parser.add_argument('--output', '-o', type=str, default=None, 
                       help='输出文件路径')
    parser.add_argument('--list-modes', action='store_true',
                       help='列出所有预设模式')
    parser.add_argument('--list-dimensions', action='store_true',
                       help='列出所有可用维度')
    
    args = parser.parse_args()
    
    # 处理特殊选项
    if args.list_modes:
        print("\n=== 预设分析模式 ===\n")
        for mode_key, mode_info in PRESET_MODES.items():
            print(f"📊 {mode_key:12} - {mode_info['name']}: {mode_info['description']}")
            print(f"   维度数量：{len(mode_info['dimensions'])}")
            print()
        sys.exit(0)
    
    if args.list_dimensions:
        print("\n=== 全部分析维度 (37 个) ===\n")
        for category, dims in DIMENSION_CATEGORIES.items():
            category_names = {
                "core": "核心九维",
                "structure": "文本结构层",
                "content": "思想内容层", 
                "knowledge": "知识体系层",
                "critical": "批判性思维层",
                "practice": "学习应用层",
                "cross": "跨文本分析层",
                "growth": "个人成长层",
                "meta": "元认知层"
            }
            print(f"\n📚 {category_names.get(category, category)}:")
            for key, (icon, desc) in dims.items():
                num = key.split('_')[0]
                print(f"   {num:2}. {key.replace('_', ' '):25} {icon}")
        sys.exit(0)
    
    # 检查文件是否存在
    if not os.path.exists(args.book_path):
        print(f"❌ 错误：文件不存在：{args.book_path}")
        sys.exit(1)
    
    # 创建分析器
    print("📖 正在提取书籍内容...")
    analyzer = EbookAnalyzer(args.book_path)
    analyzer.extract_text()
    
    if not analyzer.text_content:
        print("❌ 错误：无法从文件中提取内容，请检查格式是否支持")
        sys.exit(1)
    
    print(f"✅ 成功提取 {len(analyzer.text_content):,} 字符")
    
    # 执行命令
    if args.command == 'extract':
        print(f"\n{'='*60}")
        print(f"{analyzer.text_content[:5000]}...")  # 预览
        print(f"{'='*60}\n(总字数：{len(analyzer.text_content):,})")
        
    elif args.command == 'analyze':
        # 确定使用哪些维度
        if args.dimensions:
            # 检查是否是预设模式
            if args.dimensions.lower() in PRESET_MODES:
                mode_name = args.dimensions.lower()
                dimensions = PRESET_MODES[mode_name]['dimensions'].copy()
                if mode_name == 'all':
                    # 填充所有维度
                    dimensions = []
                    for cat_dims in DIMENSION_CATEGORIES.values():
                        dimensions.extend(list(cat_dims.keys()))
                print(f"\n📊 使用模式：{PRESET_MODES[mode_name]['name']} ({len(dimensions)}个维度)")
            else:
                # 自定义维度
                dimensions = parse_dimensions(args.dimensions)
                if not dimensions:
                    print("⚠️ 警告：未找到匹配的维度，使用默认核心九维")
                    dimensions = list(DIMENSION_CATEGORIES["core"].keys())
                print(f"\n📊 使用自定义维度 ({len(dimensions)}个): {', '.join(d[:2] for d in dimensions)}")
        else:
            # 默认使用核心九维
            dimensions = list(DIMENSION_CATEGORIES["core"].keys())
            print(f"\n📊 使用默认模式：核心九维 ({len(dimensions)}个维度)")
        
        print("🤖 正在调用本地 AI 进行分析，请稍候...")
        print(f"   → 模型：{args.model}")
        print(f"   → 维度：{len(dimensions)}个")
        print(f"   → 预计时间：{len(dimensions) * 15}秒 ~ {len(dimensions) * 60}秒\n")
        
        result = analyzer.analyze_with_ollama(args.model, dimensions)
        
        if result['success']:
            analysis = result['analysis']
            
            # 输出结果
            if args.output:
                with open(args.output, 'w', encoding='utf-8') as f:
                    f.write(analysis)
                print(f"✅ 分析完成！")
                print(f"📄 结果已保存至：{args.output}")
                print(f"\n预览前 1000 字符:\n{'='*60}\n")
                print(analysis[:1000] + "...")
            else:
                print("✅ 分析完成！\n" + "="*60)
                print(analysis)
                print("="*60)
        else:
            print(f"❌ 分析失败：{result['error']}")
            print("\n建议:")
            print("  1. 确保 Ollama 已运行：ollama serve")
            print("  2. 检查模型是否已下载：ollama pull qwen2.5")
            print("  3. 如果书籍过长，尝试减少维度数量")
            sys.exit(1)


if __name__ == '__main__':
    main()
